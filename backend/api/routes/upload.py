import os
import uuid
import datetime
import logging
import json
import httpx
from fastapi import APIRouter, UploadFile, File, Form, HTTPException, BackgroundTasks
from pydantic import BaseModel
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation

# Import custom pipelines
from vectorstore.mongodb import MongoDBVectorStore
from embeddings.bge import BGEEmbeddings
from chunking.semantic import SemanticChunker
from ocr.engine import OCREngine

logger = logging.getLogger("algonox.routes.upload")
router = APIRouter(prefix="/api/upload", tags=["upload"])

# Global state instances
vectorstore = None
embedding_client = None
semantic_chunker = None
ocr_engine = None

try:
    vectorstore = MongoDBVectorStore()
    embedding_client = BGEEmbeddings()
    semantic_chunker = SemanticChunker(embedding_client=embedding_client)
    ocr_engine = OCREngine()
except Exception as e:
    logger.error(f"Failed to initialize RAG backend instances: {e}")

def get_vectorstore():
    global vectorstore
    if not vectorstore:
        vectorstore = MongoDBVectorStore()
    return vectorstore

def get_embedding_client():
    global embedding_client
    if not embedding_client:
        embedding_client = BGEEmbeddings()
    return embedding_client

def get_semantic_chunker():
    global semantic_chunker
    if not semantic_chunker:
        semantic_chunker = SemanticChunker(embedding_client=get_embedding_client())
    return semantic_chunker

def get_ocr_engine():
    global ocr_engine
    if not ocr_engine:
        ocr_engine = OCREngine()
    return ocr_engine

def generate_document_summary(text_content: str, filename: str) -> dict:
    """
    Synchronously query Groq to generate a global summary, list of keywords,
    topics, and a semantic outline of the document.
    """
    groq_api_key = os.getenv("GROQ_API_KEY")
    if not groq_api_key:
        logger.error("GROQ_API_KEY is not configured. Skipping LLM summary generation.")
        return {
            "summary": f"Summary for {filename} is currently unavailable.",
            "keywords": ["document", "rag"],
            "topics": ["General Information"],
            "outline": [{"title": "Content Overview", "description": "Extracted text content from the uploaded file."}]
        }
    
    # Take a representative snippet of the document (up to 15,000 characters)
    sample_text = text_content[:15000]
    
    headers = {
        "Authorization": f"Bearer {groq_api_key}",
        "Content-Type": "application/json"
    }
    
    prompt = f"""You are an expert enterprise document intelligence analyzer.
Analyze the following document named "{filename}" and generate:
1. A highly professional, comprehensive executive summary (2-3 paragraphs) capturing key themes, findings, and overall purpose.
2. A list of 5-10 precise, relevant keywords/tags.
3. A list of 3-7 core topics covered.
4. A hierarchical semantic outline of key sections, chapters, or concepts discussed, each with a brief description.

You MUST respond ONLY with a valid JSON object matching this schema:
{{
  "summary": "Detailed executive summary string...",
  "keywords": ["kw1", "kw2", ...],
  "topics": ["topic1", "topic2", ...],
  "outline": [
    {{"title": "Section Title 1", "description": "Brief description..."}},
    ...
  ]
}}
Do NOT wrap the response in markdown blocks or include any other conversational text. Output ONLY strict JSON.

Document Text Snippet:
{sample_text}
"""
    
    payload = {
        "model": "llama-3.3-70b-versatile",
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.2,
        "response_format": {"type": "json_object"}
    }
    
    try:
        with httpx.Client(timeout=30.0) as client:
            response = client.post(
                "https://api.groq.com/openai/v1/chat/completions",
                headers=headers,
                json=payload
            )
            if response.status_code == 200:
                data = response.json()
                content_str = data["choices"][0]["message"]["content"].strip()
                parsed = json.loads(content_str)
                logger.info(f"Successfully generated Groq global summary for {filename}.")
                return parsed
            else:
                logger.error(f"Groq API returned status {response.status_code} during summary: {response.text}")
    except Exception as e:
        logger.error(f"Failed to generate summary using Groq: {e}")
        
    return {
        "summary": f"Executive summary of {filename}. This document contains details on various sections.",
        "keywords": [filename.split('.')[-1], "RAG"],
        "topics": ["Document Analysis"],
        "outline": [{"title": "General Content", "description": "Extracted text content from the uploaded file."}]
    }

def process_file_in_background(file_path: str, filename: str, content_type: str, document_id: str):
    """
    Asynchronous task running in the background to avoid blocking API responses.
    Parses documents, executes layout analysis, applies OCR when scanned, chunks, embeds, and saves.
    """
    try:
        page_texts = []
        active_ocr = get_ocr_engine()
        active_db = get_vectorstore()
        active_chunker = get_semantic_chunker()
        active_embedder = get_embedding_client()
        
        # 1. Parse different file formats
        if filename.endswith(".pdf"):
            logger.info(f"Ingesting native/scanned PDF: {filename}")
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                num_pages = len(reader.pages)
                for page_idx in range(num_pages):
                    page = reader.pages[page_idx]
                    page_text = page.extract_text() or ""
                    
                    # Layout Detection Fallback (verify an active OCR engine is actually initialized):
                    if len(page_text.strip()) < 50 and active_ocr:
                        active_ocr.initialize_engines()
                        if active_ocr.paddle_ocr or active_ocr.easy_ocr_reader:
                            logger.info(f"Page {page_idx + 1} of PDF '{filename}' seems scanned (low character count). Running OCR...")
                            try:
                                images_on_page = page.images
                                ocr_results = []
                                for img_idx, img_obj in enumerate(images_on_page):
                                    ocr_text = active_ocr.extract_text_from_image_bytes(img_obj.data)
                                    if ocr_text:
                                        ocr_results.append(ocr_text)
                                if ocr_results:
                                    page_text = "\n".join(ocr_results)
                            except Exception as ocr_ex:
                                logger.error(f"Failed to run inline image OCR extraction on PDF page: {ocr_ex}")
                            
                    page_texts.append((page_idx + 1, page_text))
            
        elif filename.endswith((".png", ".jpg", ".jpeg", ".bmp", ".tiff")):
            logger.info(f"Ingesting image: {filename}")
            with open(file_path, "rb") as f:
                img_bytes = f.read()
                text = active_ocr.extract_text_from_image_bytes(img_bytes)
                page_texts.append((1, text))
                
        elif filename.endswith(".docx"):
            logger.info(f"Ingesting DOCX document: {filename}")
            doc = DocxDocument(file_path)
            fullText = []
            for para in doc.paragraphs:
                fullText.append(para.text)
            page_texts.append((1, "\n".join(fullText)))
            
        elif filename.endswith(".pptx"):
            logger.info(f"Ingesting PPTX slide: {filename}")
            prs = Presentation(file_path)
            for slide_idx, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                page_texts.append((slide_idx + 1, "\n".join(slide_text)))
                
        else:
            # Fallback text file
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                page_texts.append((1, f.read()))

        # Assemble full extracted text for Global Document Summary
        full_text = "\n".join([pt[1] for pt in page_texts])

        # 2. Update uploaded file catalog metadata with correct page count (entry created during upload_files)
        active_db.files_collection.update_one(
            {"document_id": document_id},
            {"$set": {"page_count": len(page_texts) if page_texts else 1}}
        )

        # 3. Generate Global Summary & Outline
        logger.info(f"Computing LLM summary & outline for document {filename}...")
        summary_info = generate_document_summary(full_text, filename)
        summary_doc = {
            "document_id": document_id,
            "filename": filename,
            "summary": summary_info.get("summary", ""),
            "keywords": summary_info.get("keywords", []),
            "topics": summary_info.get("topics", []),
            "outline": summary_info.get("outline", []),
            "created_at": datetime.datetime.utcnow().isoformat()
        }
        active_db.insert_summary(summary_doc)

        # 4. Chunking & Indexing
        all_chunks = []
        for page_num, text_content in page_texts:
            if not text_content.strip():
                continue
                
            metadata = {
                "document_id": document_id,
                "filename": filename,
                "page_number": page_num,
                "created_at": datetime.datetime.utcnow().isoformat(),
                "file_type": content_type
            }
            
            # Use our semantic chunker
            chunks = active_chunker.chunk_document(text_content, metadata)
            all_chunks.extend(chunks)

        if not all_chunks:
            logger.warning(f"No text extracted or chunks generated for document: {filename}")
            # Update status
            active_db.files_collection.update_one({"document_id": document_id}, {"$set": {"status": "Empty"}})
            return
            
        # 5. Create Embeddings in Batches
        logger.info(f"Generating dense embeddings for {len(all_chunks)} chunks from '{filename}'")
        texts_to_embed = [c["text"] for c in all_chunks]
        embeddings = active_embedder.embed_documents(texts_to_embed)
        
        # 6. Assemble & Persist in MongoDB
        for idx, chunk in enumerate(all_chunks):
            chunk["embedding"] = embeddings[idx]
            chunk["document_id"] = document_id
            
        active_db.insert_chunks(all_chunks)
        
        # Update file catalog status to "Indexed"
        active_db.files_collection.update_one({"document_id": document_id}, {"$set": {"status": "Indexed"}})
        logger.info(f"Ingestion pipeline completed successfully for document: {filename}")
        
    except Exception as e:
        logger.error(f"Error in background document processing: {e}")
        try:
            get_vectorstore().files_collection.update_one({"document_id": document_id}, {"$set": {"status": "Error", "error": str(e)}})
        except:
            pass
    finally:
        # Cleanup temporary files
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
            except Exception as ex:
                logger.error(f"Failed deleting temporary file: {ex}")

@router.post("")
async def upload_files(
    background_tasks: BackgroundTasks,
    files: list[UploadFile] = File(...),
):
    """
    Handles multi-file simultaneous upload.
    Files are saved locally inside a temporary directory and sent to background processing queue.
    """
    uploaded_files_info = []
    
    # Establish a temp upload directory
    temp_dir = os.path.join(os.getcwd(), "temp_uploads")
    os.makedirs(temp_dir, exist_ok=True)
    
    for file in files:
        if not file.filename:
            continue
            
        # Secure Ingestion Validation
        valid_extensions = (".pdf", ".docx", ".pptx", ".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".txt")
        if not file.filename.lower().endswith(valid_extensions):
            raise HTTPException(
                status_code=400, 
                detail=f"Unsupported file extension. Supported formats: {', '.join(valid_extensions)}"
            )
            
        document_id = str(uuid.uuid4())
        file_path = os.path.join(temp_dir, f"{document_id}_{file.filename}")
        
        try:
            # Write out files in blocks
            with open(file_path, "wb") as buffer:
                while content := await file.read(1024 * 1024): # 1MB chunks
                    buffer.write(content)
            
            # Immediately record the file with "Indexing" status in the database to resolve frontend polling race conditions
            active_db = get_vectorstore()
            file_doc = {
                "_id": document_id,
                "document_id": document_id,
                "filename": file.filename,
                "file_type": file.content_type or "unknown",
                "created_at": datetime.datetime.utcnow().isoformat(),
                "page_count": 1,  # Default placeholder, will be updated in background
                "status": "Indexing"
            }
            active_db.insert_uploaded_file(file_doc)
                    
            # Launch async background ingestion thread
            background_tasks.add_task(
                process_file_in_background,
                file_path,
                file.filename,
                file.content_type or "unknown",
                document_id
            )
            
            uploaded_files_info.append({
                "document_id": document_id,
                "filename": file.filename,
                "status": "Indexing",
                "message": "File is being parsed and indexed in background."
            })
        except Exception as e:
            logger.error(f"Failed uploading file: {e}")
            raise HTTPException(status_code=500, detail=f"Failed to process upload for: {file.filename}")
            
    return {"message": "Files uploaded successfully.", "data": uploaded_files_info}

@router.get("/documents")
async def list_documents():
    """
    Retrieve unique metadata for all indexed documents.
    """
    try:
        docs = get_vectorstore().get_all_documents()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Database offline: {e}")
    # Serialize metadata for clean JSON
    result = []
    for doc in docs:
        doc_id = doc.get("document_id") or doc.get("_id")
        result.append({
            "document_id": doc_id,
            "filename": doc.get("filename", "Untitled Document"),
            "chunk_count": doc.get("chunk_count", 0),
            "created_at": doc.get("created_at", ""),
            "page_count": doc.get("page_count", 1),
            "file_type": doc.get("file_type", "unknown"),
            "status": doc.get("status", "Indexed")
        })
    return {"documents": result}

@router.delete("/documents/{document_id}")
async def delete_document(document_id: str):
    """
    Remove all document vectors and chunks from MongoDB.
    """
    try:
        count = get_vectorstore().delete_document(document_id)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Database offline: {e}")
    return {"message": f"Successfully deleted document and removed {count} chunks."}
